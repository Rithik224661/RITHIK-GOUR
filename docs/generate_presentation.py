from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()

title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Retail Insights Assistant"
subtitle.text = "System architecture, scale design and query flow"

bullet_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(bullet_slide_layout)
title = slide.shapes.title
title.text = "Architecture Overview"
body = slide.shapes.placeholders[1].text_frame
body.text = "Data ingestion and storage"
p = body.add_paragraph()
p.text = "Data pipeline: CSV ingestion -> transform -> analytics store (DuckDB/Parquet/Cloud DW)"
q = body.add_paragraph()
q.text = "LLM layer for NLU and RAG; agents handle parsing, data extraction and validation."

slide = prs.slides.add_slide(bullet_slide_layout)
title = slide.shapes.title
title.text = "LLM & Multi-agent strategy"
body = slide.shapes.placeholders[1].text_frame
body.text = "Agents"
p = body.add_paragraph()
p.text = "Language to query resolution agent"
q = body.add_paragraph()
q.text = "Data extraction agent (DuckDB, Pandas)"
r = body.add_paragraph()
r.text = "Validation agent"

slide = prs.slides.add_slide(bullet_slide_layout)
title = slide.shapes.title
title.text = "Scaling to 100GB+"
body = slide.shapes.placeholders[1].text_frame
body.text = "Storage: Parquet on S3, BigQuery or Snowflake as analytics store"
p = body.add_paragraph()
p.text = "Ingestion: Spark/Dask batch or streaming (Kafka, Firehose)"
q = body.add_paragraph()
q.text = "Retrieval: RAG + vector store (FAISS/Pinecone) for semantic filtering"

slide = prs.slides.add_slide(bullet_slide_layout)
title = slide.shapes.title
title.text = "Monitoring & Evaluation"
body = slide.shapes.placeholders[1].text_frame
body.text = "Key metrics: accuracy, latency, cost"
p = body.add_paragraph()
p.text = "Fallbacks: synthetic summaries, cached prompts, human-in-loop validation"

prs.save('docs/architecture_presentation.pptx')
print('Generated docs/architecture_presentation.pptx')
